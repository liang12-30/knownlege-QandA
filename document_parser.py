"""
文档解析模块
支持多种格式：.docx, .pdf, .xlsx, .pptx, .txt, .md, .png, .jpeg
包含板式解析和OCR功能
集成细粒度知识分块器，遵循碎片精准度>检索速度原则
"""
import os
import json
from typing import List, Dict, Any
from pathlib import Path
import hashlib

# 文档解析库
import docx
import PyPDF2
import pdfplumber
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract

# 日志
from loguru import logger

import config
from knowledge_chunker import KnowledgeChunker  # 新增：知识分块器


class DocumentParser:
    """文档解析器基类"""
    
    def __init__(self):
        """初始化解析器"""
        # 配置OCR
        if os.path.exists(config.TESSERACT_CMD):
            pytesseract.pytesseract.tesseract_cmd = config.TESSERACT_CMD
        
        self.supported_formats = ['.doc', '.docx', '.pdf', '.xlsx', '.pptx', '.txt', '.md', '.png', '.jpeg', '.jpg']
        
    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析文档
        
        Args:
            file_path: 文档路径
            
        Returns:
            解析结果字典 {doc_id, title, content, type, metadata}
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"文件不存在: {file_path}")
            return None
            
        ext = file_path.suffix.lower()
        
        if ext not in self.supported_formats:
            logger.warning(f"不支持的文件格式: {ext}")
            return None
        
        logger.info(f"开始解析文件: {file_path.name}")
        
        try:
            if ext in ['.doc', '.docx']:
                result = self._parse_word(file_path)
            elif ext == '.pdf':
                result = self._parse_pdf(file_path)
            elif ext == '.xlsx':
                result = self._parse_excel(file_path)
            elif ext == '.pptx':
                result = self._parse_pptx(file_path)
            elif ext in ['.txt', '.md']:
                result = self._parse_text(file_path)
            elif ext in ['.png', '.jpeg', '.jpg']:
                result = self._parse_image(file_path)
            else:
                return None
                
            # 生成文档ID
            doc_id = self._generate_doc_id(file_path)
            result['doc_id'] = doc_id
            result['title'] = file_path.stem
            result['type'] = ext[1:]  # 去掉点号
            
            logger.info(f"文件解析完成: {file_path.name}, 内容长度: {len(result.get('content', ''))}")
            return result
            
        except Exception as e:
            logger.error(f"解析文件失败 {file_path.name}: {str(e)}")
            return None
    
    def _generate_doc_id(self, file_path: Path) -> str:
        """生成文档唯一ID"""
        return hashlib.md5(str(file_path).encode()).hexdigest()[:16]
    
    def _parse_word(self, file_path: Path) -> Dict[str, Any]:
        """
        解析Word文档（.doc/.docx）
        
        注意：.doc格式需要转换为.docx才能解析
        """
        ext = file_path.suffix.lower()
        
        # 如果是旧版.doc格式
        if ext == '.doc':
            logger.warning(f"检测到旧版.doc格式：{file_path.name}")
            logger.warning("python-docx仅支持.docx格式")
            
            # 尝试用python-docx打开（通常会失败）
            try:
                doc = docx.Document(file_path)
            except Exception as e:
                error_msg = (
                    f"无法解析.doc文件：{file_path.name}\n"
                    f"建议：\n"
                    f"1. 使用Microsoft Word将文件另存为.docx格式\n"
                    f"2. 或使用在线转换工具将.doc转换为.docx\n"
                    f"3. 转换后重新放入dataset目录\n"
                    f"原因：{str(e)}"
                )
                logger.error(error_msg)
                return {
                    'content': f"【解析失败】{file_path.name}为旧版.doc格式，请转换为.docx后重试",
                    'tables': [],
                    'metadata': {'error': str(e), 'format': 'doc'}
                }
        else:
            # .docx格式正常解析
            doc = docx.Document(file_path)
        
        content_parts = []
        tables_data = []
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text.strip())
        
        # 提取表格
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                tables_data.append({
                    'table_id': i,
                    'data': table_data
                })
                # 将表格内容也加入文本
                content_parts.append(f"[表格{i}]")
                for row in table_data:
                    content_parts.append(" | ".join(row))
        
        return {
            'content': "\n".join(content_parts),
            'tables': tables_data,
            'metadata': {
                'paragraphs_count': len(doc.paragraphs),
                'tables_count': len(doc.tables),
                'format': ext[1:]
            }
        }
    
    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """
        解析PDF文档（使用板式解析）
        使用pdfplumber提取结构化数据
        """
        content_parts = []
        tables_data = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # 提取文本
                text = page.extract_text()
                if text:
                    content_parts.append(f"[第{page_num + 1}页]")
                    content_parts.append(text)
                
                # 提取表格（板式解析）
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    if table:
                        tables_data.append({
                            'page': page_num + 1,
                            'table_id': table_idx,
                            'data': table
                        })
                        # 将表格内容加入文本
                        content_parts.append(f"[表格-页{page_num + 1}-{table_idx}]")
                        for row in table:
                            if row:
                                content_parts.append(" | ".join([str(cell) if cell else "" for cell in row]))
        
        return {
            'content': "\n".join(content_parts),
            'tables': tables_data,
            'metadata': {
                'pages_count': len(pdf.pages) if 'pdf' in locals() else 0,
                'tables_count': len(tables_data)
            }
        }
    
    def _parse_excel(self, file_path: Path) -> Dict[str, Any]:
        """
        解析Excel文档
        提取表格数据及相邻说明文字
        """
        excel_file = pd.ExcelFile(file_path)
        
        content_parts = []
        sheets_data = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            content_parts.append(f"[工作表: {sheet_name}]")
            
            # 将DataFrame转换为文本
            # 包含列名
            if not df.empty:
                # 添加列名
                content_parts.append(" | ".join([str(col) for col in df.columns]))
                
                # 添加数据行
                for idx, row in df.iterrows():
                    row_text = " | ".join([str(val) if pd.notna(val) else "" for val in row])
                    content_parts.append(row_text)
                
                sheets_data.append({
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'columns': list(df.columns)
                })
        
        return {
            'content': "\n".join(content_parts),
            'sheets': sheets_data,
            'metadata': {
                'sheets_count': len(excel_file.sheet_names)
            }
        }
    
    def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """
        解析PPTX文档（板式解析）
        提取幻灯片内容和结构
        """
        prs = Presentation(file_path)
        
        content_parts = []
        slides_data = []
        
        for slide_idx, slide in enumerate(prs.slides):
            content_parts.append(f"[幻灯片 {slide_idx + 1}]")
            
            slide_content = []
            
            # 提取标题
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                content_parts.append(f"标题: {title_text}")
                slide_content.append({'type': 'title', 'text': title_text})
            
            # 提取其他文本框和表格
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_parts.append(shape.text.strip())
                    slide_content.append({'type': 'text', 'text': shape.text.strip()})
                
                # 提取表格
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                        content_parts.append(" | ".join(row_data))
                    
                    slide_content.append({'type': 'table', 'data': table_data})
            
            slides_data.append({
                'slide_id': slide_idx,
                'content': slide_content
            })
        
        return {
            'content': "\n".join(content_parts),
            'slides': slides_data,
            'metadata': {
                'slides_count': len(prs.slides)
            }
        }
    
    def _parse_text(self, file_path: Path) -> Dict[str, Any]:
        """解析纯文本文档"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'content': content,
            'metadata': {
                'lines_count': len(content.split('\n'))
            }
        }
    
    def _parse_image(self, file_path: Path) -> Dict[str, Any]:
        """
        解析图片文档（使用OCR）
        先转文字再解析
        """
        try:
            # 打开图片
            image = Image.open(file_path)
            
            # 使用pytesseract进行OCR识别
            # 指定中文+英文识别
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            return {
                'content': text,
                'metadata': {
                    'image_size': image.size,
                    'image_mode': image.mode,
                    'ocr_applied': True
                }
            }
        except Exception as e:
            logger.error(f"OCR识别失败 {file_path.name}: {str(e)}")
            return {
                'content': "",
                'metadata': {
                    'ocr_applied': False,
                    'error': str(e)
                }
            }


class KnowledgeBaseBuilder:
    """
    知识库构建器
    
    集成细粒度分块策略：
    - 每个文档切分为 800 字左右的片段
    - 保留 100 字重叠保持上下文
    - 按标题、表格、段落智能分块
    """
    
    def __init__(self):
        self.parser = DocumentParser()
        self.knowledge_base = []
        # 初始化知识分块器（使用配置的参数）
        self.chunker = KnowledgeChunker(
            chunk_size=config.CHUNK_SIZE,
            overlap=config.CHUNK_OVERLAP
        )
        logger.info(f"知识库构建器初始化: chunk_size={config.CHUNK_SIZE}, overlap={config.CHUNK_OVERLAP}")
    
    def build_from_directory(self, directory: str) -> List[Dict[str, Any]]:
        """
        从目录构建知识库
        
        Args:
            directory: 文档目录路径
            
        Returns:
            知识库列表
        """
        directory = Path(directory)
        
        if not directory.exists():
            logger.error(f"目录不存在: {directory}")
            return []
        
        # 获取所有支持的文件
        files = []
        for ext in self.parser.supported_formats:
            files.extend(directory.glob(f"*{ext}"))
        
        logger.info(f"找到 {len(files)} 个文档文件")
        
        # 解析所有文档并进行细粒度分块
        total_chunks = 0
        for file_path in files:
            result = self.parser.parse(str(file_path))
            if result:
                # 提取文档信息
                doc_id = result.get('doc_id', '')
                title = result.get('title', '')
                content = result.get('content', '')
                doc_type = result.get('type', '')
                
                # 使用分块器切分文档
                chunks = self.chunker.chunk_by_structure(content, title)
                
                # 将每个片段添加到知识库
                for chunk in chunks:
                    chunk_doc = {
                        'doc_id': f"{doc_id}_chunk_{chunk['chunk_index']}",
                        'title': chunk['chunk_title'],
                        'content': chunk['chunk_content'],
                        'type': doc_type,
                        # 新增字段
                        'keywords': chunk.get('keywords', []),
                        'entities': chunk.get('entities', []),
                        'chunk_type': chunk.get('chunk_type', 'paragraph'),
                        'chunk_index': chunk.get('chunk_index', 0),
                        'source_doc': title,
                        'importance_score': chunk.get('importance_score', 1.0)
                    }
                    self.knowledge_base.append(chunk_doc)
                    total_chunks += 1
                
                logger.info(f"文档 '{title}' 解析完成，切分为 {len(chunks)} 个片段")
        
        logger.info(f"知识库构建完成：{len(files)} 个文档 → {total_chunks} 个知识片段")
        return self.knowledge_base
    
    def save_to_json(self, output_path: str):
        """
        保存知识库到JSON文件
        
        Args:
            output_path: 输出文件路径
        """
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.knowledge_base, f, ensure_ascii=False, indent=2)
        
        logger.info(f"知识库已保存到: {output_path}")
    
    def load_from_json(self, input_path: str):
        """
        从JSON文件加载知识库
        
        Args:
            input_path: 输入文件路径
        """
        with open(input_path, 'r', encoding='utf-8') as f:
            self.knowledge_base = json.load(f)
        
        logger.info(f"从文件加载知识库，共 {len(self.knowledge_base)} 个文档")


if __name__ == "__main__":
    # 测试代码
    from loguru import logger
    
    logger.add("document_parser.log", rotation="10 MB")
    
    # 构建知识库
    builder = KnowledgeBaseBuilder()
    knowledge_base = builder.build_from_directory(config.KNOWLEDGE_BASE_DIR)
    
    # 保存到JSON
    output_path = os.path.join(config.PARSED_KNOWLEDGE_DIR, "knowledge_base.json")
    builder.save_to_json(output_path)
    
    print(f"知识库构建完成！共解析 {len(knowledge_base)} 个文档")
    print(f"结果已保存到: {output_path}")

